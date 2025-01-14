import PyPDF2 # type: ignore
import os
import fitz  # type: ignore
import similarity_project.app as app
import time
from googletrans import Translator # type: ignore
import pdfplumber # type: ignore
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
import pandas as pd
from reportlab.pdfgen import canvas # type: ignore
from reportlab.lib.pagesizes import letter # type: ignore
from io import BytesIO
from PIL import Image



translations = {
    "en": {
        "merge_pdfs": "Enter the PDF file names to merge (comma-separated): ",
        "output_file": "Enter the output file name: ",
        "merge_success": "Files merged successfully into: ",
        "split_pdf": "Enter the PDF file name: ",
        "start_page": "Enter the start page number (starts from 0): ",
        "end_page": "Enter the end page number: ",
        "split_success": "File split successfully into: ",
        "extract_text": "Enter the PDF file name: ",
        "text_from_page": "--- Text from page ",
        "encrypt_pdf": "Enter the PDF file name: ",
        "password": "Enter the password: ",
        "encrypt_success": "File encrypted successfully: ",
        "decrypt_pdf": "Enter the encrypted PDF file name: ",
        "decrypt_success": "File decrypted successfully: ",
        "incorrect_password": "Incorrect password!",
        "rotate_pdf": "Enter the PDF file name: ",
        "rotation_angle": "Enter the rotation angle (90, 180, 270): ",
        "rotate_success": "Pages rotated successfully in: ",
        "add_watermark": "Enter the PDF file name: ",
        "watermark_file": "Enter the watermark PDF file name: ",
        "watermark_success": "Watermark added successfully to: ",
        "compress_pdf": "Enter the PDF file name: ",
        "compress_success": "PDF compressed successfully: ",
        "extract_images": "Enter the PDF file name: ",
        "output_folder": "Enter the output folder: ",
        "extract_images_success": "Images extracted successfully to: ",
        "add_metadata": "Enter the PDF file name: ",
        "title": "Enter the title: ",
        "author": "Enter the author: ",
        "subject": "Enter the subject: ",
        "metadata_success": "Metadata added successfully to: ",
        "remove_metadata": "Enter the PDF file name: ",
        "remove_metadata_success": "Metadata removed successfully from: ",
        "rearrange_pages": "Enter the PDF file name: ",
        "page_order": "Enter the new page order (comma-separated): ",
        "rearrange_success": "Pages rearranged successfully in: ",
        "extract_pages": "Enter the PDF file name: ",
        "page_numbers": "Enter the page numbers to extract (comma-separated): ",
        "extract_pages_success": "Pages extracted successfully to: ",
        "error_occurred": "An error occurred: ",
        "function_menu": "\n--- Function Menu ---",
        "merge_files": "1. Merge PDF files",
        "split_file": "2. Split PDF file",
        "extract_text_from_pdf": "3. Extract text from PDF",
        "encrypt_file": "4. Encrypt PDF file",
        "decrypt_file": "5. Decrypt PDF file",
        "rotate_pages": "6. Rotate PDF pages",
        "add_watermark_to_pdf": "7. Add watermark to PDF",
        "compress_file": "8. Compress PDF file",
        "extract_images_from_pdf": "9. Extract images from PDF",
        "add_metadata_to_pdf": "10. Add metadata to PDF",
        "remove_metadata_from_pdf": "11. Remove metadata from PDF",
        "rearrange_pdf_pages": "12. Rearrange PDF pages",
        "extract_specific_pages": "13. Extract specific pages from PDF",
        "exit": "23. Exit",
        "enter_function_number": "Enter the function number: ",
        "invalid_choice": "Invalid choice. Please try again.",
        "goodbye": "Goodbye!"
    },
    "fr": {
        "merge_pdfs": "Entrez les noms des fichiers PDF à fusionner (séparés par des virgules) : ",
        "output_file": "Entrez le nom du fichier de sortie : ",
        "merge_success": "Fichiers fusionnés avec succès dans : ",
        "split_pdf": "Entrez le nom du fichier PDF : ",
        "start_page": "Entrez le numéro de la page de début (commence à 0) : ",
        "end_page": "Entrez le numéro de la page de fin : ",
        "split_success": "Fichier divisé avec succès dans : ",
        "extract_text": "Entrez le nom du fichier PDF : ",
        "text_from_page": "--- Texte de la page ",
        "encrypt_pdf": "Entrez le nom du fichier PDF : ",
        "password": "Entrez le mot de passe : ",
        "encrypt_success": "Fichier crypté avec succès : ",
        "decrypt_pdf": "Entrez le nom du fichier PDF crypté : ",
        "decrypt_success": "Fichier décrypté avec succès : ",
        "incorrect_password": "Mot de passe incorrect !",
        "rotate_pdf": "Entrez le nom du fichier PDF : ",
        "rotation_angle": "Entrez l'angle de rotation (90, 180, 270) : ",
        "rotate_success": "Pages tournées avec succès dans : ",
        "add_watermark": "Entrez le nom du fichier PDF : ",
        "watermark_file": "Entrez le nom du fichier PDF de filigrane : ",
        "watermark_success": "Filigrane ajouté avec succès à : ",
        "compress_pdf": "Entrez le nom du fichier PDF : ",
        "compress_success": "PDF compressé avec succès : ",
        "extract_images": "Entrez le nom du fichier PDF : ",
        "output_folder": "Entrez le dossier de sortie : ",
        "extract_images_success": "Images extraites avec succès dans : ",
        "add_metadata": "Entrez le nom du fichier PDF : ",
        "title": "Entrez le titre : ",
        "author": "Entrez l'auteur : ",
        "subject": "Entrez le sujet : ",
        "metadata_success": "Métadonnées ajoutées avec succès à : ",
        "remove_metadata": "Entrez le nom du fichier PDF : ",
        "remove_metadata_success": "Métadonnées supprimées avec succès de : ",
        "rearrange_pages": "Entrez le nom du fichier PDF : ",
        "page_order": "Entrez le nouvel ordre des pages (séparés par des virgules) : ",
        "rearrange_success": "Pages réorganisées avec succès dans : ",
        "extract_pages": "Entrez le nom du fichier PDF : ",
        "page_numbers": "Entrez les numéros de pages à extraire (séparés par des virgules) : ",
        "extract_pages_success": "Pages extraites avec succès dans : ",
        "error_occurred": "Une erreur s'est produite : ",
        "function_menu": "\n--- Menu des fonctions ---",
        "merge_files": "1. Fusionner des fichiers PDF",
        "split_file": "2. Diviser un fichier PDF",
        "extract_text_from_pdf": "3. Extraire du texte d'un PDF",
        "encrypt_file": "4. Crypter un fichier PDF",
        "decrypt_file": "5. Décrypter un fichier PDF",
        "rotate_pages": "6. Tourner les pages d'un PDF",
        "add_watermark_to_pdf": "7. Ajouter un filigrane à un PDF",
        "compress_file": "8. Compresser un fichier PDF",
        "extract_images_from_pdf": "9. Extraire des images d'un PDF",
        "add_metadata_to_pdf": "10. Ajouter des métadonnées à un PDF",
        "remove_metadata_from_pdf": "11. Supprimer les métadonnées d'un PDF",
        "rearrange_pdf_pages": "12. Réorganiser les pages d'un PDF",
        "extract_specific_pages": "13. Extraire des pages spécifiques d'un PDF",
        "exit": "23. Quitter",
        "enter_function_number": "Entrez le numéro de la fonction : ",
        "invalid_choice": "Choix invalide. Veuillez réessayer.",
        "goodbye": "Au revoir !"
    },
        "ar": {
        "merge_pdfs": "أدخل أسماء ملفات PDF للفصل (مفصولة بفواصل): ",
        "output_file": "أدخل اسم الملف الناتج: ",
        "merge_success": "تم دمج الملفات بنجاح في: ",
        "split_pdf": "أدخل اسم ملف PDF: ",
        "start_page": "أدخل رقم الصفحة الأولى (يبدأ من 0): ",
        "end_page": "أدخل رقم الصفحة الأخيرة: ",
        "split_success": "تم تقسيم الملف إلى: ",
        "extract_text": "أدخل اسم ملف PDF: ",
        "text_from_page": "--- نص الصفحة ",
        "encrypt_pdf": "أدخل اسم ملف PDF: ",
        "password": "أدخل كلمة المرور: ",
        "encrypt_success": "تم تشفير الملف: ",
        "decrypt_pdf": "أدخل اسم ملف PDF المشفر: ",
        "decrypt_success": "تم فك تشفير الملف: ",
        "incorrect_password": "كلمة المرور غير صحيحة!",
        "rotate_pdf": "أدخل اسم ملف PDF: ",
        "rotation_angle": "أدخل زاوية الدوران (90، 180، 270): ",
        "rotate_success": "تم تدوير الصفحات بنجاح في: ",
        "add_watermark": "أدخل اسم ملف PDF: ",
        "watermark_file": "أدخل اسم ملف PDF العلامة المائية: ",
        "watermark_success": "تمت إضافة العلامة المائية بنجاح إلى: ",
        "compress_pdf": "أدخل اسم ملف PDF: ",
        "compress_success": "تم ضغط ملف PDF بنجاح: ",
        "extract_images": "أدخل اسم ملف PDF: ",
        "output_folder": "أدخل المجلد الناتج: ",
        "extract_images_success": "تم استخراج الصور بنجاح إلى: ",
        "add_metadata": "أدخل اسم ملف PDF: ",
        "title": "أدخل العنوان: ",
        "author": "أدخل المؤلف: ",
        "subject": "أدخل الموضوع: ",
        "metadata_success": "تمت إضافة البيانات الوصفية بنجاح إلى: ",
        "remove_metadata": "أدخل اسم ملف PDF: ",
        "remove_metadata_success": "تمت إزالة البيانات الوصفية بنجاح من: ",
        "rearrange_pages": "أدخل اسم ملف PDF: ",
        "page_order": "أدخل ترتيب الصفحات الجديد (مفصولة بفواصل): ",
        "rearrange_success": "تمت إعادة ترتيب الصفحات بنجاح في: ",
        "extract_pages": "أدخل اسم ملف PDF: ",
        "page_numbers": "أدخل أرقام الصفحات لاستخراجها (مفصولة بفواصل): ",
        "extract_pages_success": "تم استخراج الصفحات بنجاح إلى: ",
        "error_occurred": "حدث خطأ: ",
         "function_menu": "\n--- قائمة الوظائف ---",
        "merge_files": "1. دمج ملفات PDF",
        "split_file": "2. تقسيم ملف PDF",
        "extract_text_from_pdf": "3. استخراج النصوص من PDF",
        "encrypt_file": "4. تشفير ملف PDF",
        "decrypt_file": "5. فك تشفير ملف PDF",
        "rotate_pages": "6. تدوير صفحات PDF",
        "add_watermark_to_pdf": "7. إضافة علامة مائية إلى PDF",
         "compress_file": "8. ضغط ملف PDF",
        "extract_images_from_pdf": "9. استخراج الصور من PDF",
        "add_metadata_to_pdf": "10. إضافة بيانات وصفية إلى PDF",
        "remove_metadata_from_pdf": "11. إزالة البيانات الوصفية من PDF",
         "rearrange_pdf_pages": "12. إعادة ترتيب صفحات PDF",
        "extract_specific_pages": "13. استخراج صفحات محددة من PDF",
        "exit": "23. خروج",
        "enter_function_number": "أدخل رقم الوظيفة: ",
        "invalid_choice": "اختيار غير صحيح. حاول مرة أخرى.",
         "goodbye": "وداعًا!"
    }
}



def get_translation(lang, key):
    return translations[lang].get(key, key)

def get_static_input_files():
    return ["1.pdf", "2.pdf"]


def merge_pdfs(lang):
    print("mergeing 1.pdf and 2.pdf file from the server !")
    pdf_list = get_static_input_files()
    output_file = "output_file" + ".pdf"
    merger = PdfMerger()
    try:
        for pdf in pdf_list:
            pdf = pdf.strip()
            if pdf:
                merger.append(pdf)
        merger.write(output_file)
        print(get_translation(lang, "merge_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))
    finally:
        merger.close()


def split_pdf(lang):
    print("spliting 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    try:
        start = int(input(get_translation(lang, "start_page")))
        end = int(input(get_translation(lang, "end_page")))
        output_file = "output_file" + ".pdf"
        reader = PdfReader(input_file)
        writer = PdfWriter()
        for page in range(start, end):
            writer.add_page(reader.pages[page])
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "split_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def extract_text(lang):
    print("extract text from 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".txt"
    reader = PdfReader(input_file)
    try:
        with open(output_file, "w", encoding="utf-8") as txt_file:
            for page_num, page in enumerate(reader.pages):
                txt_file.write(get_translation(lang, "text_from_page") + str(page_num + 1) + " ---\n")
                txt_file.write(page.extract_text() + "\n")
        print(get_translation(lang, "extract_pages_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def encrypt_pdf(lang):
    print("encrypting 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    password = input(get_translation(lang, "password"))
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        writer.append(reader)
        writer.encrypt(password)
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "encrypt_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def decrypt_pdf(lang):
    print("decrypting 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    password = input(get_translation(lang, "password"))
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        if reader.decrypt(password):
            writer.append(reader)
            with open(output_file, "wb") as output_pdf:
                writer.write(output_pdf)
            print(get_translation(lang, "decrypt_success") + output_file)
        else:
            print(get_translation(lang, "incorrect_password"))
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def rotate_pdf(lang):
    print("rotateing 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    rotation = int(input(get_translation(lang, "rotation_angle")))
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        for page in reader.pages:
            page.rotate(rotation)
            writer.add_page(page)
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "rotate_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def add_watermark(lang):
    print("adding watermark to 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    print("adding watermark to 2.pdf file from the server !")
    watermark_file = get_static_input_files()[1]
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    watermark = PdfReader(watermark_file)
    writer = PdfWriter()
    try:
        watermark_page = watermark.pages[0]
        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "watermark_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def compress_pdf(lang):
    print("compressing 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        for page in reader.pages:
            writer.add_page(page)
        writer.add_metadata(reader.metadata)
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "compress_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def extract_images(lang):
    print("Extracting images from 1.pdf file from the server!")
    input_file = get_static_input_files()[0]
    output_folder = "images_1_pdf"
    
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    pdf_document = fitz.open(input_file)
    
    try:
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            images = page.get_images(full=True)

            for img_index, img in enumerate(images):
                xref = img[0]
                image = pdf_document.extract_image(xref)
                image_data = image["image"]
                image_filename = f"{output_folder}/image_{page_num + 1}_{img_index + 1}.jpg"
                
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_data)
                
        print(get_translation(lang, "extract_images_success") + output_folder)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))
    finally:
        pdf_document.close()


def add_metadata(lang):
    print("adding metadata to 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".pdf"
    title = input(get_translation(lang, "title"))
    author = input(get_translation(lang, "author"))
    subject = input(get_translation(lang, "subject"))
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        writer.append(reader)
        writer.add_metadata({
            "/Title": title,
            "/Author": author,
            "/Subject": subject
        })
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "metadata_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))

def view_metadata(lang):
    try:
        reader = PdfReader("output_file.pdf")
        metadata = reader.metadata
        print("Metadata for file:", "output_file.pdf")
        print("Title:", metadata.get("/Title", "No Title"))
        print("Author:", metadata.get("/Author", "No Author"))
        print("Subject:", metadata.get("/Subject", "No Subject"))
    except Exception as e:
        print(f"Error occurred while reading metadata: {str(e)}")

def remove_metadata(lang):
    print("removing metadata from 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".pdf"
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        writer.append(reader)
        writer.add_metadata({})
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "remove_metadata_success") + output_file)
    except Exception as e:
         print(get_translation(lang, "error_occurred") + str(e))


def rearrange_pages(lang):
    print("rearranging pages of 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".pdf"
    page_order = input(get_translation(lang, "page_order")).split(',')
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        for page_num in page_order:
            writer.add_page(reader.pages[int(page_num.strip())])
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "rearrange_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def extract_pages(lang):
    print("extracting pages of 1.pdf file from the server !")
    input_file = get_static_input_files()[0]
    output_file = "output_file" + ".pdf"
    page_numbers = input(get_translation(lang, "page_numbers")).split(',')
    reader = PdfReader(input_file)
    writer = PdfWriter()
    try:
        for page_num in page_numbers:
            writer.add_page(reader.pages[int(page_num.strip())])
        with open(output_file, "wb") as output_pdf:
            writer.write(output_pdf)
        print(get_translation(lang, "extract_pages_success") + output_file)
    except Exception as e:
        print(get_translation(lang, "error_occurred") + str(e))


def pdf_doce():
    print("Converting PDF to Word document")
    try:
        # Create a Word document object
        doc = Document()

        # Open the PDF using pdfplumber
        with pdfplumber.open("1.pdf") as pdf:
            # Iterate through each page of the PDF
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()

                if text:  # Ensure there is text on the page
                    doc.add_heading(f'Page {page_num + 1}', level=1)
                    # Add the extracted text as a paragraph
                    doc.add_paragraph(text)

        # Save the Word document
        doc.save("output.docx")
        print(f"PDF successfully converted to Word document at output.docx")

    except Exception as e:
        print(f"An error occurred: {e}")


def pdf_excel():
    print("Converting PDF to Excel")
    try:
        # Open the PDF file using pdfplumber
        with pdfplumber.open("1.pdf") as pdf:
            all_tables = []
            # Iterate through all the pages of the PDF
            for page in pdf.pages:
                # Extract tables from the page
                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table[1:], columns=table[0])  # Convert table to DataFrame
                    all_tables.append(df)
            
            # Concatenate all tables into a single DataFrame (if multiple tables exist)
            full_data = pd.concat(all_tables, ignore_index=True)

            # Save the data to an Excel file
            full_data.to_excel("output.xlsx", index=False, engine='openpyxl')
            print(f"PDF data successfully saved to output_presentation.pptx")
    
    except Exception as e:
        print(f"An error occurred: {e}")


def pdf_html():
    print("Converting PDF to HTML")
    try:
        # Open the PDF using pdfplumber
        with pdfplumber.open("1.pdf") as pdf:
            html_content = "<html><body>"
            
            # Iterate through each page of the PDF
            for page_num, page in enumerate(pdf.pages):
                # Extract text from the page
                text = page.extract_text()
                
                if text:  # If the page has text
                    # Add a header for each page
                    html_content += f"<h2>Page {page_num + 1}</h2>"
                    # Convert the extracted text into paragraphs
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        html_content += f"<p>{para}</p>"
            
            # Closing HTML tags
            html_content += "</body></html>"
            
            # Save the HTML content to a file
            with open("output.html", 'w', encoding='utf-8') as file:
                file.write(html_content)
            print(f"PDF successfully converted to HTML at output.html")

    except Exception as e:
        print(f"An error occurred: {e}")


def pdf_ppt():
    print("Converting PDF to PowerPoint")
    try:
        # Create a PowerPoint presentation object
        prs = Presentation()

        # Open the PDF using pdfplumber
        with pdfplumber.open("1.pdf") as pdf:
            # Iterate through each page of the PDF
            for page_num, page in enumerate(pdf.pages):
                # Extract text from the page
                text = page.extract_text()

                if text:  # Ensure there is text on the page
                    # Add a slide to the presentation for each page of the PDF
                    slide_layout = prs.slide_layouts[1]  # Use a title and content layout
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    content = slide.shapes.placeholders[1]

                    # Set the title to the page number and content to the extracted text
                    title.text = f"Page {page_num + 1}"
                    content.text = text

        # Save the PowerPoint presentation
        prs.save("output_presentation.pptx")
        print("PDF successfully converted to PowerPoint at output_presentation.pptx")

    except Exception as e:
        print(f"An error occurred: {e}")


def extract_text_from_pdf(pdf_path):
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
    except Exception as e:
        return f"An error occurred while extracting text: {e}"

def translate_text(text, target_language='ar'):
    try:
        translator = Translator()
        translated = translator.translate(text, dest=target_language)
        return translated.text
    except Exception as e:
        return f"An error occurred during translation: {e}"

def save_text_to_file(text, filename='translated_text.txt'):
    try:
        with open(filename, 'w', encoding='utf-8') as file:
            file.write(text)
        print(f"Translated text saved to {filename}")
    except Exception as e:
        print(f"An error occurred while saving the file: {e}")

def transalte():
    print("Translating text from 1.pdf to Arabic")
    pdf_path = '1.pdf'  # Replace with your PDF file path
    target_lang = 'ar'  # For Arabic translation

    extracted_text = extract_text_from_pdf(pdf_path)
    print(extracted_text)
    if "An error occurred" not in extracted_text:
        translated_text = translate_text(extracted_text, target_lang)
        save_text_to_file(translated_text)

def clear_screen():
    # يعمل على Windows
    time.sleep(1.5)  
    os.system('cls')



def add_page_numbers():
    print("Adding page numbers to the 1.pdf file from the server!")
    try:
        # Open the original PDF
        with open("1.pdf", 'rb') as original_pdf:
            reader = PyPDF2.PdfReader(original_pdf) # type: ignore
            writer = PyPDF2.PdfWriter()

            # Create a temporary PDF to store the page numbers
            packet = BytesIO()
            for page_num in range(len(reader.pages)):
                # Create a canvas for each page
                c = canvas.Canvas(packet)
                c.drawString(500, 10, f"Page {page_num + 1}")  # Position the page number at the bottom right
                c.save()

                # Go back to the beginning of the BytesIO buffer
                packet.seek(0)

                # Create a new PDF with the page number
                page_number_pdf = PyPDF2.PdfReader(packet)
                page = reader.pages[page_num]
                page_number = page_number_pdf.pages[0]

                # Merge the page with the page number
                page.merge_page(page_number)

                # Add the merged page to the writer
                writer.add_page(page)

            # Write the output to a new PDF file
            with open("output.pdf", 'wb') as output_pdf:
                writer.write(output_pdf)

            print(f"Page numbers added successfully! The new PDF is saved as output.pdf")

    except Exception as e:
        print(f"An error occurred: {e}")



def signture():
    print("Adding signature as 1.png to the 1.pdf file from the server!")
    x_position = 450  # X-coordinate for the signature
    y_position = 30   # Y-coordinate for the signature

    try:
        # Open the original PDF
        with open("1.pdf", 'rb') as original_pdf:
            reader = PyPDF2.PdfReader(original_pdf)
            writer = PyPDF2.PdfWriter()

            # Load the signature image
            signature_image = "1.png"

            for page_num, page in enumerate(reader.pages):
                # Create a BytesIO buffer for the overlay
                packet = BytesIO()
                c = canvas.Canvas(packet, pagesize=letter)

                # Draw the signature image
                c.drawImage(signature_image, x_position, y_position, width=100, height=50)
                c.save()

                # Move to the beginning of the BytesIO buffer
                packet.seek(0)

                # Create a new PDF with the signature overlay
                overlay_pdf = PyPDF2.PdfReader(packet)
                overlay_page = overlay_pdf.pages[0]

                # Merge the overlay with the current page
                page.merge_page(overlay_page)

                # Add the merged page to the writer
                writer.add_page(page)

            # Write the output to a new PDF file
            with open("output.pdf", 'wb') as output_pdf:
                writer.write(output_pdf)

            print("Signature added successfully! The new PDF is saved as 'output.pdf'.")

    except Exception as e:
        print(f"An error occurred: {e}")


def main():
    clear_screen()
    print("Welcome to the PDF Manipulation and Similarity Script CLI!")
    lang = input("Choose language (en/fr/ar): ").strip().lower()
    if lang not in translations:
        print("Invalid language choice. Defaulting to English.")
        lang = "en"

    while True:
        clear_screen()
        print(get_translation(lang, "function_menu"))
        print(get_translation(lang, "merge_files"))
        print(get_translation(lang, "split_file"))
        print(get_translation(lang, "extract_text_from_pdf"))
        print(get_translation(lang, "encrypt_file"))
        print(get_translation(lang, "decrypt_file"))
        print(get_translation(lang, "rotate_pages"))
        print(get_translation(lang, "add_watermark_to_pdf"))
        print(get_translation(lang, "compress_file"))
        print(get_translation(lang, "extract_images_from_pdf"))
        print(get_translation(lang, "add_metadata_to_pdf"))
        print(get_translation(lang, "remove_metadata_from_pdf"))
        print(get_translation(lang, "rearrange_pdf_pages"))
        print(get_translation(lang, "extract_specific_pages"))
        print("14. View metadata")
        print("15. View Simalrity")
        print("16. PDF to EXCEL")
        print("17. PDF to DOC")
        print("18. PDF to PPT")
        print("19. PDF to HTML")
        print("20. Add page numbers")
        print("21. Signture")
        print("22. Translate to arabic")
        print(get_translation(lang, "exit"))
        choice = input(get_translation(lang, "enter_function_number"))
        
        if choice == "1":
            clear_screen()
            merge_pdfs(lang)

        elif choice == "2":
            clear_screen()
            split_pdf(lang)

        elif choice == "3":
            clear_screen()
            extract_text(lang)

        elif choice == "4":
            clear_screen()
            encrypt_pdf(lang)

        elif choice == "5":
            clear_screen()
            decrypt_pdf(lang)

        elif choice == "6":
            clear_screen()
            rotate_pdf(lang)

        elif choice == "7":
            clear_screen()
            add_watermark(lang)

        elif choice == "8":
             clear_screen()
             compress_pdf(lang)

        elif choice == "9":
            clear_screen()
            extract_images(lang)

        elif choice == "10":
            clear_screen()
            add_metadata(lang)

        elif choice == "11":
            clear_screen()
            remove_metadata(lang)

        elif choice == "12":
            clear_screen()
            rearrange_pages(lang)

        elif choice == "13":
            clear_screen()
            extract_pages(lang)

        elif choice == "14":
            clear_screen()
            view_metadata(lang)

        elif choice == "15":
            clear_screen()
            app.main()

        elif choice == "16":
            clear_screen()
            pdf_excel()

        elif choice == "17":
            clear_screen()
            pdf_doce()

        elif choice == "18":
            clear_screen()
            pdf_ppt()

        elif choice == "19":
            clear_screen()
            pdf_html()

        elif choice == "20":
            clear_screen()
            add_page_numbers()

        elif choice == "21":
            clear_screen()
            signture()

        elif choice == "22":
            clear_screen()
            transalte()

        elif choice == "23":
            clear_screen()
            print(get_translation(lang, "goodbye"))
            clear_screen()
            break

        else:
            print(get_translation(lang, "invalid_choice"))

if __name__ == "__main__":
    main()


